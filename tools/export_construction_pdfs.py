import argparse
import os
from pathlib import Path
from typing import Iterable


def iter_pdfs(input_dir: Path) -> Iterable[Path]:
    for path in sorted(input_dir.rglob("*.pdf")):
        if path.is_file():
            yield path


def safe_stem(path: Path) -> str:
    # Keep Chinese names; just avoid trailing dots/spaces.
    stem = path.stem.strip().strip(".")
    return stem or "unnamed"


def export_pdf_to_png(pdf_path: Path, out_dir: Path, zoom: float) -> list[Path]:
    import fitz  # PyMuPDF

    doc = fitz.open(pdf_path)
    exported: list[Path] = []

    pdf_out_dir = out_dir / safe_stem(pdf_path)
    pdf_out_dir.mkdir(parents=True, exist_ok=True)

    matrix = fitz.Matrix(zoom, zoom)

    for index in range(doc.page_count):
        page = doc.load_page(index)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        out_path = pdf_out_dir / f"page-{index + 1:03d}.png"
        pix.save(out_path.as_posix())
        exported.append(out_path)

    return exported


def build_pptx(pdf_name: str, image_paths: list[Path], out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Use wide layout; scale images to fit slide.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)

        # Fit image into slide while keeping aspect ratio
        from PIL import Image

        with Image.open(image_path) as img:
            img_w, img_h = img.size

        slide_w = prs.slide_width
        slide_h = prs.slide_height

        img_ratio = img_w / img_h
        slide_ratio = slide_w / slide_h

        if img_ratio >= slide_ratio:
            # fit to width
            width = slide_w
            height = int(slide_w / img_ratio)
            left = 0
            top = int((slide_h - height) / 2)
        else:
            # fit to height
            height = slide_h
            width = int(slide_h * img_ratio)
            top = 0
            left = int((slide_w - width) / 2)

        slide.shapes.add_picture(image_path.as_posix(), left, top, width=width, height=height)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path.as_posix())


def main() -> int:
    parser = argparse.ArgumentParser(description="Export construction PDFs into per-page PNGs + PPTX for easy annotation.")
    parser.add_argument("--input", default=str(Path(__file__).resolve().parents[1] / "施工图纸"), help="Input folder containing PDFs")
    parser.add_argument(
        "--output",
        default=str(Path(__file__).resolve().parents[1] / "施工图纸" / "导出_可编辑"),
        help="Output folder",
    )
    parser.add_argument("--zoom", type=float, default=2.5, help="Render zoom (2.0~3.0 is typical)")
    parser.add_argument("--no-pptx", action="store_true", help="Only export PNGs")

    args = parser.parse_args()

    input_dir = Path(args.input)
    out_dir = Path(args.output)

    if not input_dir.exists():
        raise SystemExit(f"Input dir not found: {input_dir}")

    out_dir.mkdir(parents=True, exist_ok=True)

    pdfs = list(iter_pdfs(input_dir))
    if not pdfs:
        print(f"No PDFs found under: {input_dir}")
        return 0

    index_lines = [
        "# 施工图纸导出清单\n",
        f"- 输入目录: {input_dir}\n",
        f"- 输出目录: {out_dir}\n",
        f"- 渲染倍率(zoom): {args.zoom}\n",
        "\n## 文件列表\n",
    ]

    any_pptx_failed = False

    for pdf_path in pdfs:
        print(f"[export] {pdf_path.name}")
        image_paths = export_pdf_to_png(pdf_path, out_dir, args.zoom)

        rel_folder = (out_dir / safe_stem(pdf_path)).relative_to(out_dir)
        index_lines.append(f"- {pdf_path.name} → {rel_folder.as_posix()}/ (共 {len(image_paths)} 页 PNG)\n")

        if not args.no_pptx:
            pptx_dir = out_dir / "PPT"
            pptx_path = pptx_dir / f"{safe_stem(pdf_path)}.pptx"
            try:
                build_pptx(pdf_path.name, image_paths, pptx_path)
                index_lines.append(f"  - PPT: PPT/{pptx_path.name}\n")
            except Exception as e:
                any_pptx_failed = True
                index_lines.append(f"  - PPT: 生成失败（{type(e).__name__}: {e}）\n")

    (out_dir / "INDEX.md").write_text("".join(index_lines), encoding="utf-8")

    print(f"\nDone. See: {out_dir / 'INDEX.md'}")
    if any_pptx_failed:
        print("Note: Some PPTX exports failed. You can re-run with --no-pptx to only export PNGs.")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
